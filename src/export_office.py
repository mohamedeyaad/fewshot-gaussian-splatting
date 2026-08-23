"""Export the report to .docx and the deck to .pptx.

  python src/export_office.py

The HTML versions stay canonical - they are what the build scripts generate and
what every number is read into. These are conversions for submission, produced
from those same files rather than written by hand, so they cannot drift.

  results/report_brief.html -> results/report.docx
  results/report.html       -> results/report_full.docx
  results/slides.html       -> results/presentation.pptx

Images are embedded as data URIs in the source HTML; they are decoded back to
bytes here rather than re-read from disk, so the exports match exactly what the
HTML shows.
"""
from __future__ import annotations

import base64
import io
import os
import re
from pathlib import Path

from bs4 import BeautifulSoup, NavigableString, Tag

import docx
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor

from pptx import Presentation
from pptx.dml.color import RGBColor as PRGB
from pptx.util import Emu, Inches as PInches, Pt as PPt

ROOT = Path(os.path.expanduser("~/fewshot_gs"))
RES = ROOT / "results"

INK = RGBColor(0x16, 0x20, 0x2A)
MUTED = RGBColor(0x5C, 0x6E, 0x78)
ACCENT = RGBColor(0x0B, 0x6E, 0x7F)
GAIN = RGBColor(0x1C, 0x7C, 0x54)
LOSS = RGBColor(0xB3, 0x45, 0x2C)

P_INK = PRGB(0x16, 0x20, 0x2A)
P_MUTED = PRGB(0x5C, 0x6E, 0x78)
P_ACCENT = PRGB(0x0B, 0x6E, 0x7F)
P_GAIN = PRGB(0x1C, 0x7C, 0x54)
P_LOSS = PRGB(0xB3, 0x45, 0x2C)


def soup_of(path: Path) -> BeautifulSoup:
    return BeautifulSoup(path.read_text(encoding="utf-8"), "html.parser")


def img_bytes(tag: Tag):
    """Decode a data: URI back to bytes. Returns None for anything else -
    the CSP means there should be nothing else, but a missing figure should
    degrade to a skipped image rather than a traceback."""
    src = tag.get("src", "")
    m = re.match(r"data:image/[^;]+;base64,(.*)$", src, re.S)
    if not m:
        return None
    return io.BytesIO(base64.b64decode(m.group(1)))


def signed_colour(text: str):
    """Deltas carry their meaning in the sign, and the HTML colours them. Keep
    that: a table of signed numbers in one colour loses the finding."""
    t = text.strip()
    if re.match(r"^\+\d", t):
        return GAIN, P_GAIN
    if re.match(r"^[-−]\d", t):
        return LOSS, P_LOSS
    return None, None


# ---------------------------------------------------------------- DOCX ----

def add_runs(par, node, base_size=10.5):
    """Walk inline children so bold/italic/code survive the conversion."""
    for child in node.children:
        if isinstance(child, NavigableString):
            txt = str(child)
            if txt.strip() or txt == " ":
                r = par.add_run(re.sub(r"\s+", " ", txt))
                r.font.size = Pt(base_size)
                r.font.color.rgb = INK
            continue
        if not isinstance(child, Tag):
            continue
        txt = re.sub(r"\s+", " ", child.get_text())
        if not txt.strip():
            continue
        r = par.add_run(txt)
        r.font.size = Pt(base_size)
        r.font.color.rgb = INK
        if child.name in ("b", "strong"):
            r.bold = True
        elif child.name in ("i", "em"):
            r.italic = True
        elif child.name == "code":
            r.font.name = "Consolas"
            r.font.size = Pt(base_size - 1)
            r.font.color.rgb = ACCENT


def docx_table(doc, table: Tag):
    head = table.find("thead")
    body = table.find("tbody")
    if body is None:
        return
    head_rows = head.find_all("tr") if head else []
    body_rows = body.find_all("tr")
    ncol = max((len(r.find_all(["td", "th"])) for r in head_rows + body_rows),
               default=0)
    if not ncol:
        return

    t = doc.add_table(rows=0, cols=ncol)
    t.style = "Table Grid"
    t.alignment = WD_TABLE_ALIGNMENT.CENTER

    for tr in head_rows:
        cells = tr.find_all(["th", "td"])
        row = t.add_row().cells
        for i, c in enumerate(cells[:ncol]):
            p = row[i].paragraphs[0]
            r = p.add_run(re.sub(r"\s+", " ", c.get_text()).strip())
            r.bold = True
            r.font.size = Pt(8.5)
            r.font.color.rgb = MUTED

    for tr in body_rows:
        cells = tr.find_all(["td", "th"])
        row = t.add_row().cells
        for i, c in enumerate(cells[:ncol]):
            txt = re.sub(r"\s+", " ", c.get_text()).strip()
            p = row[i].paragraphs[0]
            r = p.add_run(txt)
            r.font.size = Pt(9)
            col, _ = signed_colour(txt)
            r.font.color.rgb = col or INK
            if c.find(["b", "strong"]):
                r.bold = True
            if i > 0:
                p.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    cap = table.find("caption")
    if cap:
        p = doc.add_paragraph()
        r = p.add_run(re.sub(r"\s+", " ", cap.get_text()).strip())
        r.font.size = Pt(8.5)
        r.italic = True
        r.font.color.rgb = MUTED


def html_to_docx(src: Path, dst: Path, title: str):
    soup = soup_of(src)
    doc = docx.Document()
    for s in doc.sections:
        s.top_margin = s.bottom_margin = Inches(0.9)
        s.left_margin = s.right_margin = Inches(0.9)
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(10.5)

    h = doc.add_paragraph()
    r = h.add_run(title)
    r.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = INK
    sub = doc.add_paragraph()
    r = sub.add_run("Mohamed Eyad  ·  Università degli Studi di Genova")
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED

    seen = set()
    for node in soup.find_all(["h1", "h2", "h3", "p", "table", "img", "ul", "ol"]):
        if any(a is node or node in getattr(a, "descendants", []) for a in seen):
            continue
        if node.find_parent("table") or node.find_parent(["ul", "ol"]):
            continue
        if node.name in ("h1", "h2", "h3"):
            txt = re.sub(r"\s+", " ", node.get_text()).strip()
            # The HTML numbers sections with a <span class="n">; keep the text.
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(14 if node.name != "h3" else 10)
            r = p.add_run(txt)
            r.bold = True
            r.font.size = Pt({"h1": 16, "h2": 14, "h3": 11.5}[node.name])
            r.font.color.rgb = INK if node.name != "h3" else ACCENT
        elif node.name == "p":
            txt = node.get_text().strip()
            if not txt:
                continue
            p = doc.add_paragraph()
            p.paragraph_format.space_after = Pt(6)
            add_runs(p, node)
        elif node.name in ("ul", "ol"):
            for li in node.find_all("li", recursive=False):
                p = doc.add_paragraph(style="List Bullet")
                add_runs(p, li)
            seen.add(node)
        elif node.name == "table":
            docx_table(doc, node)
            seen.add(node)
        elif node.name == "img":
            b = img_bytes(node)
            if b:
                doc.add_picture(b, width=Inches(6.4))
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
                fig = node.find_parent("figure")
                cap = fig.find("figcaption") if fig else None
                if cap:
                    p = doc.add_paragraph()
                    r = p.add_run(re.sub(r"\s+", " ", cap.get_text()).strip())
                    r.font.size = Pt(8.5)
                    r.italic = True
                    r.font.color.rgb = MUTED

    doc.save(str(dst))
    print(f"wrote {dst}  ({dst.stat().st_size/1024:,.0f} KB)")


# ---------------------------------------------------------------- PPTX ----

def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def draw_bars(slide, rows, left, top, width):
    """Redraw the diverging bar chart as PowerPoint shapes.

    The chart is CSS divs in the HTML, so a text-only conversion silently drops
    it - and it sits on the slide carrying the headline finding, which would
    have arrived in the exam as a title and no data. The geometry is already
    computed in the markup (width:N% of the track, class pos/neg for the side),
    so it is reused rather than recalculated.
    """
    from pptx.enum.shapes import MSO_SHAPE

    lab_w = PInches(1.85)
    val_w = PInches(1.15)
    track_l = left + lab_w + PInches(0.15)
    track_w = width - lab_w - val_w - PInches(0.30)
    centre = track_l + track_w / 2
    row_h = PInches(0.46)
    y = top

    for row in rows:
        lab = row.select_one(".barlab")
        val = row.select_one(".barval")
        bar = row.select_one(".bar")

        tf = textbox(slide, left, y, lab_w, row_h)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = lab.get_text().strip() if lab else ""
        r.font.size = PPt(14)
        r.font.color.rgb = P_MUTED

        track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, track_l,
                                       y + PInches(0.07), track_w,
                                       PInches(0.30))
        track.fill.solid()
        track.fill.fore_color.rgb = PRGB(0xF2, 0xF5, 0xF6)
        track.line.color.rgb = PRGB(0xD2, 0xDD, 0xE1)
        track.shadow.inherit = False

        zero = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, centre,
                                      y + PInches(0.03), Emu(9144),
                                      PInches(0.38))
        zero.fill.solid()
        zero.fill.fore_color.rgb = P_MUTED
        zero.line.fill.background()
        zero.shadow.inherit = False

        if bar is not None:
            m = re.search(r"width:\s*([\d.]+)%", bar.get("style", ""))
            frac = float(m.group(1)) / 100.0 if m else 0.0
            bw = max(int(track_w * frac), 3000)
            pos = "pos" in (bar.get("class") or [])
            bl = centre if pos else centre - bw
            shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bl,
                                         y + PInches(0.09), bw, PInches(0.26))
            shp.fill.solid()
            shp.fill.fore_color.rgb = P_GAIN if pos else P_LOSS
            shp.line.fill.background()
            shp.shadow.inherit = False

        tf = textbox(slide, track_l + track_w + PInches(0.15), y, val_w, row_h)
        p = tf.paragraphs[0]
        p.alignment = 2                       # right
        r = p.add_run()
        txt = val.get_text().strip() if val else ""
        r.text = txt
        r.font.size = PPt(14)
        r.font.bold = True
        _, pcol = signed_colour(txt)
        r.font.color.rgb = pcol or P_INK

        y += row_h
    return y


def pptx_from_slides(src: Path, dst: Path):
    soup = soup_of(src)
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    M = PInches(0.72)
    W = prs.slide_width - 2 * M

    for sec in soup.find_all("section", class_="slide"):
        slide = prs.slides.add_slide(blank)
        y = PInches(0.55)

        eyebrow = sec.find(class_="eyebrow")
        if eyebrow:
            tf = textbox(slide, M, y, W, PInches(0.32))
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = re.sub(r"\s+", " ", eyebrow.get_text()).strip().upper()
            r.font.size = PPt(11)
            r.font.bold = True
            r.font.color.rgb = P_ACCENT
            y += PInches(0.42)

        head = sec.find(["h1", "h2"])
        if head:
            txt = re.sub(r"\s+", " ", head.get_text()).strip()
            big = head.name == "h1"
            tf = textbox(slide, M, y, W, PInches(1.5 if big else 1.0))
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = txt
            r.font.size = PPt(34 if big else 26)
            r.font.bold = True
            r.font.color.rgb = P_INK
            y += PInches(1.25 if big else 0.95)
            if big and len(txt) > 60:
                y += PInches(0.35)

        barrows = sec.select(".barrow")
        if barrows:
            y = draw_bars(slide, barrows, M, y, W) + PInches(0.2)

        img = sec.find("img")
        table = sec.find("table")
        body = [p for p in sec.find_all("p")
                if p.get_text().strip() and not p.find_parent("figure")]
        bullets = sec.find_all("li")
        cards = sec.select(".card")

        if img is not None:
            b = img_bytes(img)
            if b:
                avail_h = prs.slide_height - y - PInches(0.7)
                pic = slide.shapes.add_picture(b, M, y, width=W)
                if pic.height > avail_h:            # never overflow the slide
                    scale = avail_h / pic.height
                    pic.height = int(pic.height * scale)
                    pic.width = int(pic.width * scale)
                    pic.left = int(M + (W - pic.width) / 2)
            continue

        if table is not None:
            rows_h = table.find("thead").find_all("tr") if table.find("thead") else []
            rows_b = table.find("tbody").find_all("tr") if table.find("tbody") else []
            allr = rows_h + rows_b
            ncol = max((len(r.find_all(["td", "th"])) for r in allr), default=0)
            if ncol and allr:
                h = PInches(min(0.42 * len(allr) + 0.15, 3.4))
                gt = slide.shapes.add_table(len(allr), ncol, M, y, W, h).table
                for ri, tr in enumerate(allr):
                    cells = tr.find_all(["td", "th"])
                    for ci, c in enumerate(cells[:ncol]):
                        txt = re.sub(r"\s+", " ", c.get_text()).strip()
                        cell = gt.cell(ri, ci)
                        cell.text = txt
                        para = cell.text_frame.paragraphs[0]
                        for run in para.runs:
                            run.font.size = PPt(15)
                            _, pcol = signed_colour(txt)
                            run.font.color.rgb = pcol or P_INK
                            if ri < len(rows_h):
                                run.font.bold = True
                y += h + PInches(0.25)

        if cards:
            n = len(cards)
            cw = int((W - PInches(0.25) * (n - 1)) / n)
            for i, card in enumerate(cards):
                k = card.select_one(".k")
                cap = card.select_one(".cap")
                h3 = card.find("h3")
                tf = textbox(slide, M + i * (cw + PInches(0.25)), y,
                             cw, PInches(2.0))
                if k:
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = k.get_text().strip()
                    r.font.size = PPt(30)
                    r.font.bold = True
                    r.font.color.rgb = P_INK
                if h3:
                    p = tf.add_paragraph()
                    r = p.add_run()
                    r.text = h3.get_text().strip()
                    r.font.size = PPt(15)
                    r.font.bold = True
                    r.font.color.rgb = P_INK
                if cap:
                    p = tf.add_paragraph()
                    r = p.add_run()
                    r.text = re.sub(r"\s+", " ", cap.get_text()).strip()
                    r.font.size = PPt(12)
                    r.font.color.rgb = P_MUTED
            y += PInches(2.15)

        remaining = prs.slide_height - y - PInches(0.5)
        if remaining > PInches(0.6) and (body or bullets):
            tf = textbox(slide, M, y, W, remaining)
            first = True
            for li in bullets:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                r = p.add_run()
                r.text = "•  " + re.sub(r"\s+", " ", li.get_text()).strip()
                r.font.size = PPt(15)
                r.font.color.rgb = P_INK
                p.space_after = PPt(9)
            if not bullets:
                for pr in body[:3]:
                    if pr.find_parent(".card"):
                        continue
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    r = p.add_run()
                    r.text = re.sub(r"\s+", " ", pr.get_text()).strip()
                    r.font.size = PPt(16 if "lead" in (pr.get("class") or []) else 14)
                    r.font.color.rgb = (P_MUTED if "muted" in (pr.get("class") or [])
                                        else P_INK)
                    p.space_after = PPt(10)

    prs.save(str(dst))
    print(f"wrote {dst}  ({dst.stat().st_size/1024:,.0f} KB)")


if __name__ == "__main__":
    html_to_docx(RES / "report_brief.html", RES / "report.docx",
                 "Few-Shot Gaussian Splatting with Diffusion-Based Augmentation")
    html_to_docx(RES / "report.html", RES / "report_full.docx",
                 "Diffusion Augmentation for Few-Shot Splatting — full report")
    pptx_from_slides(RES / "slides.html", RES / "presentation.pptx")
