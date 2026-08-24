"""The spoken script, injected into the deck and exported as a document.

  python src/speaker_notes.py

Writes the notes into results/presentation.pptx (PowerPoint presenter view
shows them) and into results/speaker_notes.docx for reading on paper.

The notes are matched to slides by the slide TITLE, not by index, so inserting
or reordering a slide in build_slides.py does not silently shift every note by
one. A slide with no matching note, or a note matching no slide, is reported
rather than passed over - that mismatch is exactly what would go unnoticed
until the exam.

Written to be spoken, not read out. Short sentences, one idea each, and the
numbers that matter said aloud. Timings assume ~13 minutes of talking.
"""
from __future__ import annotations

import os
import re
from pathlib import Path

import docx
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation

ROOT = Path(os.path.expanduser("~/fewshot_gs"))
RES = ROOT / "results"

INK = RGBColor(0x16, 0x20, 0x2A)
MUTED = RGBColor(0x5C, 0x6E, 0x78)
ACCENT = RGBColor(0x0B, 0x6E, 0x7F)

# (title fragment, seconds, [paragraphs])
NOTES = [
 ("Where this goes", 22, [
  "Quick map of the next fifteen minutes.",
  "The problem, then the method and the pipeline, then the grid I ran. Then "
  "the result, which is a sign change. Then why it happens, and the controls "
  "that rule out the alternatives. A prediction I tested. And the limits.",
  "The one thing to carry through: the same treatment can help or harm, and "
  "which one it does depends on how many real photographs you already have.",
 ]),

 ("When does a generated photograph", 45, [
  "Good morning. My project asks one question: can AI-generated images "
  "replace photographs you never took, when you are reconstructing a 3D scene?",
  "The short answer is yes, but only in a narrow window — and I spent most of "
  "the project working out where that window ends and why.",
  "236 training runs, two scenes, and every number here is read straight out "
  "of the raw results files, not typed in by hand.",
 ]),

 ("3D Gaussian Splatting is excellent", 60, [
  "Gaussian Splatting reconstructs a 3D scene from photographs. Given enough "
  "of them it is excellent — 25.2 dB on this outdoor scene from 219 photos.",
  "Cut it to five photos and it collapses to 15.2. That is a ten decibel gap.",
  "Photographs are expensive to collect. A diffusion model generates images "
  "for free. So the question is whether free images can buy back that gap.",
 ]),

 ("what it looks like", 40, [
  "This is the same held-out camera reconstructed from 5, 10, 20 and 219 photos.",
  "The important thing is what fails at five views. The shape of the truck is "
  "broadly right. What costs the ten decibels is the haze, the floaters, and "
  "the surfaces the optimiser had no second view to pin down.",
  "So the deficit is coverage, not detail. That matters later.",
 ]),

 ("Three ways to manufacture", 60, [
  "I tested three strategies.",
  "Inpainting masks part of a real photo and refills it. Same camera, same "
  "framing — no new geometry at all. It is the control, really.",
  "Outpainting keeps the lens and enlarges the frame. The camera does not "
  "move; it just sees wider. That supplies peripheral content.",
  "Pose-guided is the ambitious one. Estimate depth, warp the image to a new "
  "camera position I choose, and let diffusion fill the holes. It is the only "
  "one that invents a genuinely new viewpoint.",
  "Each is generated at four ratios, three subset sizes, three seeds.",
 ]),

 ("One condition, end to end", 48, [
  "This is one condition, start to finish.",
  "Two hundred and nineteen photographs go through COLMAP for camera poses and "
  "sparse points. Then the split, and this is the important box: thirty-two "
  "photographs are held out and never trained on, in any run, verified by hash.",
  "From the training pool I select k views by farthest-point sampling — five, "
  "ten or twenty.",
  "Then the fork. The top lane trains on those k real images. The bottom lane "
  "sends them through diffusion augmentation first, and trains on k real plus "
  "n synthetic.",
  "Both lanes hit the same trainer for the same seven thousand iterations and "
  "are scored on the same thirty-two held-out photographs. The only thing that "
  "differs between the arms is what went into training.",
 ]),

 ("What was actually run", 38, [
  "And this is the grid.",
  "Three strategies, four synthetic ratios, three subset sizes, three seeds. "
  "A hundred and eight augmented runs, plus baselines, controls and a "
  "full-data ceiling.",
  "Each tick in a cell is one training run.",
  "The critical detail is the pairing. Every augmented cell is compared "
  "against the baseline built from the SAME seed's subset — so the luck of "
  "which five views got drawn cancels out, instead of being counted as an "
  "effect of augmentation.",
 ]),

 ("survive a 0.2 dB effect", 55, [
  "The effects here are small, so the measurement has to be careful.",
  "Every run of a scene is scored against byte-identical held-out images — I "
  "verify that by hash across all 236 runs.",
  "Every augmented run is compared against its own seed's baseline. Which five "
  "photos you happen to draw matters enormously, and pairing cancels it.",
  "And I measured the noise floor: the identical configuration run three times "
  "varies by 0.039 dB. So a 0.2 dB effect is real and a 0.03 dB effect is not.",
 ]),

 ("helps at five views and harms at twenty", 75, [
  "This is the main result.",
  "Outpainting at the 200 percent ratio. At five real views it gains 0.285 dB. "
  "At ten it is zero. At twenty it LOSES 0.618.",
  "The treatment is identical. The only thing that changed is how many real "
  "photographs it was added to.",
  "PAUSE HERE. Let them look at the sign change.",
  "For scale: five more real photographs are worth 1.9 dB. Augmentation buys "
  "you about a seventh of one photograph, and only while photographs are scarce.",
 ]),

 ("Result &middot; what it looks like", 40, [
  "The same thing visually. Left pair is five views, right pair is twenty.",
  "One caution on this figure: per-view effects scatter several dB either side "
  "of the average. So I picked the views closest to the mean effect and printed "
  "the per-view number on each. Three arbitrary views can contradict the finding "
  "they are supposed to illustrate.",
 ]),

 ("Every cell, in one picture", 52, [
  "Here is the whole grid filled in. Green helps, red harms.",
  "Read the middle row — outpainting — left to right across the three panels. "
  "Green at five views. Neutral at ten. Red at twenty.",
  "That is the finding. Same treatment, same ratio, opposite sign.",
  "The top row, inpainting, is flat everywhere — it reuses the camera pose "
  "exactly. The bottom row, pose-guided, is red everywhere and gets darker as "
  "k grows — it invents a whole new camera.",
  "Colour is clipped at plus or minus one decibel, otherwise pose-guided's "
  "magnitude washes the crossover out.",
 ]),

 ("It reproduces indoors", 45, [
  "The obvious worry is that this is one scene's quirk. So I repeated the "
  "sweep on an indoor room from a different dataset — different capture "
  "geometry, different scale.",
  "Same sign flip. And indoors the effect is stronger at five views.",
  "Note the third number: the baselines scatter 0.52 dB between seeds indoors, "
  "against 0.07 outdoors. Which five views you draw matters far more in a room. "
  "Without the paired design I could not have resolved the effect at all.",
 ]),

 ("Two terms, one decaying", 65, [
  "Here is why the sign flips.",
  "Every synthetic view gives you two things at once. Coverage — some part of "
  "the scene no real camera saw. And inconsistency — it is invented, so it "
  "disagrees with the real photographs.",
  "Coverage DECAYS as you add real views. You have already seen that area.",
  "Inconsistency does NOT decay. A contradiction is just as harmful at twenty "
  "views as at five. Arguably worse, because now it contradicts a "
  "better-determined geometry.",
  "A shrinking benefit plus a constant cost gives you a sign change. That is "
  "the whole mechanism.",
  "But that is just a story. So I made it predict things that could fail.",
 ]),

 ("A second checkpoint reverses", 55, [
  "First objection: maybe this is a quirk of Stable Diffusion 1.5.",
  "So I swapped in Dreamshaper-8 — same architecture, same memory budget, but "
  "tuned for better photorealism.",
  "It reverses in the same place. Positive at five, negative at twenty.",
  "And here is the interesting part: the better-looking model is WORSE at every "
  "subset size. It even crosses over earlier.",
  "A model tuned to make each image individually more convincing has no reason "
  "to be more consistent BETWEEN images. Which is the point — what a synthetic "
  "view contributes is coverage, not photorealism.",
 ]),

 ("Constraint without a viewpoint", 92, [
  "Second and sharper prediction.",
  "If the harm comes from inventing a camera, then something that helps the "
  "geometry WITHOUT inventing a camera should never turn negative.",
  "Depth regularisation is exactly that. A depth network predicts a depth map "
  "for each real photograph, anchored to true scale against the sparse "
  "reconstruction points. No camera is invented, no pixel is fabricated.",
  "It is positive at every subset size — including twenty, where outpainting "
  "loses 0.6.",
  "Coverage crosses over. Constraint does not. The two differ in exactly one "
  "respect, and it is the one the theory says matters.",
  "And they compound: together at five views they reach 0.714 dB, the best "
  "result in the study — because they repair different deficiencies.",
  "One more thing, and it is the strongest evidence I have. Train to thirty "
  "thousand iterations and the outpainting gain dies — minus 0.078, nothing. "
  "The depth prior at the same length is still plus 0.208, and the two together "
  "plus 0.469. Fabricated pixels lose their value to longer training. A "
  "geometric constraint does not.",
 ]),

 ("What both scenes say together", 50, [
  "So, four things.",
  "Augmentation is not free coverage. It helps only in the single digits.",
  "The mechanism is pose novelty — four controls rule out image quality, hole "
  "content, warping artifacts, and mere repetition.",
  "Constraint beats invention.",
  "And the two are separable, which is why combining them beats either alone.",
  "Practically: below about ten real views, use it. Above that, go and take "
  "more photographs.",
 ]),

 ("What this does not show", 88, [
  "The honest limits.",
  "The biggest one is training length. I train for 7,000 iterations. I measured "
  "three points — the gain decays from plus 0.285, to plus 0.045 at fifteen "
  "thousand, to minus 0.078 at thirty thousand.",
  "But — and this is the point — the plain baselines do not improve either. Flat "
  "from seven to fifteen thousand, and lower by thirty. With five photos, long "
  "training memorises them instead of learning the scene. So 30,000 is not "
  "better converged, it is over-trained, and 7,000 is a fair operating point. "
  "The augmentation gain is still conditional on early stopping — the depth "
  "prior is not.",
  "Two things run against my own claim. At twenty views the prior is only plus "
  "0.130 after long training, inside the noise. And the compounding is "
  "single-scene: on the indoor scene the two are additive, not super-additive.",
  "Three seeds is a consistency check, not a significance test. Both diffusion "
  "checkpoints share an architecture, because SDXL and FLUX do not fit in four "
  "gigabytes.",
 ]),

 ("Augmentation buys coverage", 35, [
  "To close.",
  "Augmentation buys coverage and pays for it in consistency. Whether that "
  "trade is worth taking depends entirely on how scarce your real photographs "
  "are.",
  "At five views, worth it. By twenty, the same treatment is actively harmful.",
  "Thank you — happy to take questions.",
 ]),

 ("Thank you", 12, [
  "Thank you — happy to take questions.",
 ]),

]

QA = [
 ("Why 7,000 iterations and not 30,000?",
  "Because 30,000 is past the optimum for few-shot. I measured three lengths: "
  "the plain 5-view baseline is 15.20 at seven thousand, 15.22 at fifteen "
  "thousand — flat, inside the noise floor — and 15.04 at thirty. With few "
  "photos, longer training overfits. I state in the report that the "
  "augmentation gain is conditional on early stopping. The depth prior is the "
  "exception: it still returns plus 0.208 at thirty thousand."),
 ("You chose the five views optimally with farthest-point sampling. Isn't that cheating?",
  "It is the best case, so I also trained the random draws. Augmentation does "
  "not rescue badly-chosen views — it is no better there, and worse at twenty. "
  "Outpainting widens the frustum around cameras you already have; it cannot "
  "reach scene a badly-spread set never came near."),
 ("Is this just a Stable Diffusion artifact?",
  "No. A second checkpoint reverses in the same place. And the more "
  "photorealistic one is worse everywhere, which is evidence the mechanism is "
  "coverage rather than image quality."),
 ("Only three seeds — is any of this significant?",
  "Three seeds supports a consistency check, not a formal significance test, "
  "and I say so. The asterisk means the mean exceeds the between-seed spread. "
  "What carries the weight is that the effect reproduces across two scenes, two "
  "checkpoints, and agrees across three metrics with different failure modes."),
 ("Why is pose-guided worst, when it supplies the most new information?",
  "Because it invents the most. It is the only strategy that fabricates a "
  "camera that never existed, so it carries the most inconsistency. The theory "
  "predicts it should be worst, and it is — that is a point in favour of the "
  "account, not against it."),
 ("How do you know the numbers are right?",
  "Four automated checks over every run: that all runs of a scene are scored "
  "against byte-identical held-out images, that train and test never overlap, "
  "that subset sizes are what the manifest claims, and that baselines improve "
  "monotonically with more real views. They caught two real bugs during the "
  "project."),
 ("What would you do with more compute?",
  "A multi-view-consistent generator — Zero123++ or SV3D. The whole failure "
  "mode here is that a 2D model has no reason to be consistent between views. "
  "Those models are consistent by construction. They need more than four "
  "gigabytes, which is why this project could not use them."),
]


def norm(s):
    return re.sub(r"\s+", " ", s).strip().lower()


def main():
    pptx_path = RES / "presentation.pptx"
    prs = Presentation(str(pptx_path))

    titles = []
    for slide in prs.slides:
        txt = " ".join(sh.text_frame.text for sh in slide.shapes
                       if sh.has_text_frame)
        titles.append(norm(txt))

    used, total = set(), 0
    for frag, secs, paras in NOTES:
        f = norm(re.sub(r"&middot;", "·", frag))
        hit = next((i for i, t in enumerate(titles)
                    if f in t and i not in used), None)
        if hit is None:
            print(f"  !! no slide matches note: {frag!r}")
            continue
        used.add(hit)
        total += secs
        tf = prs.slides[hit].notes_slide.notes_text_frame
        tf.text = f"[{secs}s]"
        for p in paras:
            para = tf.add_paragraph()
            para.text = p
    for i, t in enumerate(titles):
        if i not in used:
            print(f"  !! slide {i+1} has no note: {t[:56]}")

    prs.save(str(pptx_path))
    print(f"notes injected into {pptx_path.name}  "
          f"({len(used)}/{len(titles)} slides, {total//60}m {total%60}s)")

    # ---- the readable version -----------------------------------------
    doc = docx.Document()
    for s in doc.sections:
        s.top_margin = s.bottom_margin = Inches(0.8)
        s.left_margin = s.right_margin = Inches(0.9)
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(11)

    p = doc.add_paragraph()
    r = p.add_run("Presentation script")
    r.bold = True
    r.font.size = Pt(20)
    r.font.color.rgb = INK
    p = doc.add_paragraph()
    r = p.add_run(f"Few-Shot Gaussian Splatting  ·  {len(used)} slides  ·  "
                  f"about {total//60} minutes of speaking")
    r.font.size = Pt(10)
    r.font.color.rgb = MUTED

    for i, (frag, secs, paras) in enumerate(NOTES, 1):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(15)
        r = p.add_run(f"Slide {i}")
        r.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = INK
        r = p.add_run(f"   ({secs}s)")
        r.font.size = Pt(10)
        r.font.color.rgb = ACCENT
        for para in paras:
            q = doc.add_paragraph()
            q.paragraph_format.space_after = Pt(5)
            emph = para.isupper() or para.startswith("PAUSE")
            r = q.add_run(para)
            r.font.size = Pt(11)
            r.italic = emph
            r.font.color.rgb = ACCENT if emph else INK

    doc.add_page_break()
    p = doc.add_paragraph()
    r = p.add_run("Likely questions")
    r.bold = True
    r.font.size = Pt(16)
    r.font.color.rgb = INK
    for q, a in QA:
        pq = doc.add_paragraph()
        pq.paragraph_format.space_before = Pt(11)
        r = pq.add_run(q)
        r.bold = True
        r.font.size = Pt(11)
        r.font.color.rgb = INK
        pa = doc.add_paragraph()
        r = pa.add_run(a)
        r.font.size = Pt(11)
        r.font.color.rgb = INK

    out = RES / "speaker_notes.docx"
    doc.save(str(out))
    print(f"wrote {out}  ({out.stat().st_size/1024:,.0f} KB)")


def cheat_sheet():
    """One page to hold while presenting.

    Not a summary of the report - a recall aid. Only what is hard to
    reconstruct under pressure: the exact numbers, the mechanism in three
    lines, and the one-sentence answer to each likely question.
    """
    doc = docx.Document()
    for s in doc.sections:
        s.top_margin = s.bottom_margin = Inches(0.5)
        s.left_margin = s.right_margin = Inches(0.55)
    doc.styles["Normal"].font.name = "Calibri"
    doc.styles["Normal"].font.size = Pt(9)

    def head(txt, size=11, colour=ACCENT, before=7):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(before)
        p.paragraph_format.space_after = Pt(2)
        r = p.add_run(txt.upper())
        r.bold = True
        r.font.size = Pt(size)
        r.font.color.rgb = colour
        return p

    def line(txt, bold=False, size=9):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(1)
        r = p.add_run(txt)
        r.bold = bold
        r.font.size = Pt(size)
        r.font.color.rgb = INK
        return p

    p = doc.add_paragraph()
    r = p.add_run("Few-Shot Gaussian Splatting — one page")
    r.bold = True
    r.font.size = Pt(15)
    r.font.color.rgb = INK

    head("The three numbers", 11, ACCENT, 4)
    line("5 real views  →  +0.285 dB     10 →  −0.003     20 →  −0.618", True, 10)
    line("Same treatment (outpainting, 200%). Only the number of real photos changed.")
    line("Scale: 5 more REAL photos = +1.91 dB. Ceiling (219 views) = 25.23. Floor (5) = 15.20.")

    head("The mechanism, in three lines")
    line("A fake view gives COVERAGE + CONTRADICTIONS, together.")
    line("Coverage decays as real views accumulate. Contradictions do not.")
    line("Shrinking benefit + constant cost = sign change.")

    head("Predictions it made")
    line("Depth prior (constraint, no invented camera) → never crosses over.  "
         "+0.259 / +0.163 / +0.155 at k=5/10/20.  HELD", True)
    line("Combined with outpainting at k=5: +0.714 dB — best result in the study.")
    line("At 30k: outpaint −0.078 (dead), depth +0.208, both +0.469 — constraint outlasts "
         "invention.")
    line("Against me: k=20 depth at 30k is +0.130 (n.s.); drjohnson interaction +0.075 "
         "± 0.672 — additive, not super-additive.")
    line("Pose-guided invents most → should hurt most. −1.013 at k=5.  HELD", True)

    head("Controls (what they rule out)")
    line("Duplicate real views  → not mere view count")
    line("Warp-only, no diffusion → −3.94 vs −1.45: diffusion REPAIRS, the pose is the damage")
    line("Dreamshaper-8 → reverses in the same place; better-looking model is WORSE everywhere")
    line("Noise floor → σ = 0.039 dB, so 0.2 dB is real")

    head("Second scene (drjohnson, indoor)")
    line("Same sign flip. Stronger at k=5. Baselines scatter 0.52 dB between seeds vs 0.07 "
         "outdoors — which is why pairing matters.")

    head("The honest limits")
    line("Training length: gain decays +0.285 → +0.045 (15k) → −0.078 (30k). Baseline "
         "flat 7k→15k (15.20 → 15.22), lower at 30k (15.04) — few-shot overfits, so "
         "early stopping is correct. Depth prior survives 30k: +0.208.", True)
    line("Random view selection: augmentation does NOT rescue badly-chosen views.")
    line("Three seeds = consistency check, not significance. Both checkpoints are SD 1.5-class.")

    head("If asked…")
    for q, a in QA:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(1)
        r = p.add_run(q + "  ")
        r.bold = True
        r.font.size = Pt(8.5)
        r.font.color.rgb = INK
        r = p.add_run(re.sub(r"\s+", " ", a))
        r.font.size = Pt(8.5)
        r.font.color.rgb = MUTED

    out = RES / "cheat_sheet.docx"
    doc.save(str(out))
    print(f"wrote {out}  ({out.stat().st_size/1024:,.0f} KB)")


if __name__ == "__main__":
    main()
    cheat_sheet()
